from pptx import Presentation
from pptx.util import Inches, Pt

from ..base import BimbBase
from .image import BimbImgProcessor

class BimbPPTX(BimbBase):
  def __init__(self, *args, **kwargs):
    super().__init__(**kwargs)
    self.PPTX_Obj = Presentation()
    self.image_per_row = 4
    self.image_per_col = 4

    self.image_spacing = 0.25

    self.image_width = Inches((self.PPTX_Obj.slide_width.inches / self.image_per_col) - self.image_spacing)
    self.image_height = Inches((self.PPTX_Obj.slide_height.inches / self.image_per_row) - self.image_spacing)

    self.row_index = 0
    self.col_index = 0

    self.layout = self.PPTX_Obj.slide_masters[0].slide_layouts[0]
    self.slide = self.PPTX_Obj.slides.add_slide(self.layout)

  def set_images(self, _image_details_df, img_cols=['img_hbonds', 'img_full']):
    _Bimp = BimbImgProcessor(self.path_base)

    for _idx, _imgdet in _image_details_df.iterrows():
      for _img in img_cols:
        _img_jpeg = _imgdet[_img]
        _img = self.change_ext(_img_jpeg, 'png')
        _Bimp.clip_remove_bg(_img_jpeg, _img)
        _left = (self.image_width * self.col_index)  + Inches(self.image_spacing)
        _top = (self.image_height * self.row_index) + Inches(self.image_spacing)

        # set height or width (preserve aspect ratio) or both (stretch)

        # _pic_ph = self.slide.shapes.add_picture(_img, _left, _top, self.image_width, self.image_height)
        _pic_ph = self.slide.shapes.add_picture(_img, _left, _top, None, self.image_height) # preserve aspect ratio

        # Add text with text properties (Text and Font-size)

        txBox = self.slide.shapes.add_textbox(_left, _top - Inches(0.2), self.image_width, self.image_height/4)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = f"{_imgdet.Receptor__PDB_ID_} {_imgdet.Ligand}"
        p.font.size = Pt(12)

        txBox = self.slide.shapes.add_textbox(_left, _top + Inches(0.2), self.image_width, self.image_height/4)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = f"{_imgdet.VINA_Score} kcal/mol"
        p.font.size = Pt(12)

        if self.col_index < self.image_per_col -1:
          self.col_index += 1
        else:
          self.row_index += 1
          self.col_index = 0

        if self.row_index < self.image_per_row:
          ...
        else:
          self.slide = self.PPTX_Obj.slides.add_slide(self.layout)
          self.row_index = 0

  def save(self, path):
    self.PPTX_Obj.save(path)

# _ppx = ArrangePPTXImages()
# _image_cols = ['img_contacts', 'img_full']
# _ppx.set_images(_excel_df.sort_values(['Receptor__PDB_ID_', 'Ligand']), _image_cols)
# _ppx.save(DMan.get_path("Docking-Result-Images-Tiled--Contacts.pptx"))
